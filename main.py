import os
import re
import json
import sqlite3
from datetime import datetime
from typing import Optional, List

import requests
from bs4 import BeautifulSoup
from dotenv import load_dotenv

from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse

from pydantic import BaseModel, Field
from openai import OpenAI

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


# ============================================================
# CONFIG
# ============================================================

load_dotenv()

API_KEY = os.getenv("OPENAI_API_KEY")

if not API_KEY:
    raise RuntimeError(
        "OPENAI_API_KEY is missing from .env"
    )

client = OpenAI(api_key=API_KEY)

app = FastAPI(
    title="VC Brain",
    version="5.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"]
)

DATABASE = "vc_brain.db"
UPLOAD_DIR = "uploads"

os.makedirs(UPLOAD_DIR, exist_ok=True)

MAX_URL_TEXT = 50000
MAX_DOCUMENT_TEXT = 100000
MAX_TOTAL_EVIDENCE = 200000


# ============================================================
# 300+ FACTORS
# ============================================================

FACTOR_GROUPS = {

"problem": [
"problem_severity",
"problem_frequency",
"problem_urgency",
"problem_cost",
"problem_revenue_impact",
"problem_risk_impact",
"problem_emotional_intensity",
"problem_inefficiency",
"problem_frustration",
"problem_unmet_need",
"problem_customer_awareness",
"problem_willingness_to_solve",
"problem_painkiller_strength",
"problem_mission_criticality",
"problem_persistence",
"problem_growth",
"problem_timing",
"problem_structural_nature",
"problem_regulatory_pressure",
"problem_technology_enablement"
],

"customer": [
"customer_clarity",
"ideal_customer_profile",
"customer_segmentation",
"customer_size",
"customer_accessibility",
"customer_affordability",
"customer_willingness_to_pay",
"customer_ability_to_pay",
"buyer_identity",
"user_identity",
"decision_maker_identity",
"economic_buyer_identity",
"customer_concentration",
"customer_diversity",
"customer_retention",
"customer_loyalty",
"customer_satisfaction",
"customer_dissatisfaction",
"customer_feedback_quality",
"customer_reference_quality",
"customer_testimonial_quality",
"customer_repeat_usage",
"customer_frequency",
"customer_lifetime_potential",
"customer_expansion_potential",
"customer_referral_potential",
"customer_churn_risk",
"customer_switching_cost",
"customer_migration_cost",
"customer_contract_lockin",
"customer_trust",
"customer_behavior_change"
],

"market": [
"tam",
"sam",
"som",
"tam_realism",
"sam_realism",
"som_realism",
"market_growth",
"market_growth_acceleration",
"market_size",
"market_timing",
"market_maturity",
"market_fragmentation",
"market_concentration",
"market_structure",
"market_cycle",
"market_cyclicality",
"market_resilience",
"market_stability",
"market_volatility",
"market_geographic_expansion",
"market_international_potential",
"market_adjacent_expansion",
"market_category_creation",
"market_winner_take_most",
"market_winner_take_all",
"market_network_effect_potential",
"market_billion_dollar_potential",
"market_trillion_dollar_potential",
"market_demand_growth",
"market_supply_constraints",
"market_demand_supply_gap",
"industry_growth",
"industry_profitability",
"industry_spending",
"industry_consolidation",
"industry_disruption_potential",
"industry_technology_change",
"industry_regulatory_change",
"industry_demographic_change",
"industry_behavior_change"
],

"product": [
"product_quality",
"product_usefulness",
"product_reliability",
"product_performance",
"product_speed",
"product_accuracy",
"product_security",
"product_privacy",
"product_safety",
"product_scalability",
"product_availability",
"product_uptime",
"product_usability",
"product_design",
"product_user_experience",
"product_customer_experience",
"product_accessibility",
"product_personalization",
"product_customization",
"product_integrations",
"product_compatibility",
"product_interoperability",
"product_roadmap",
"product_roadmap_quality",
"product_iteration_speed",
"product_release_frequency",
"product_feedback_loop",
"product_adoption",
"product_activation",
"product_engagement",
"product_retention",
"product_expansion",
"product_stickiness",
"product_virality",
"product_network_effect"
],

"innovation": [
"technical_innovation",
"scientific_innovation",
"product_innovation",
"business_model_innovation",
"distribution_innovation",
"pricing_innovation",
"manufacturing_innovation",
"design_innovation",
"workflow_innovation",
"process_innovation",
"data_innovation",
"ai_innovation",
"software_innovation",
"hardware_innovation",
"research_innovation",
"innovation_importance",
"innovation_customer_value",
"innovation_difficulty",
"innovation_defensibility",
"innovation_replicability",
"innovation_timing",
"innovation_market_relevance",
"innovation_breakthrough_level"
],

"differentiation": [
"differentiation",
"differentiation_clarity",
"differentiation_importance",
"differentiation_customer_value",
"differentiation_durability",
"differentiation_defensibility",
"differentiation_copy_difficulty",
"differentiation_price_advantage",
"differentiation_quality_advantage",
"differentiation_speed_advantage",
"differentiation_convenience_advantage",
"differentiation_distribution_advantage",
"differentiation_data_advantage",
"differentiation_brand_advantage",
"differentiation_technology_advantage"
],

"pricing": [
"price_level",
"price_vs_competitors",
"price_vs_customer_value",
"price_elasticity",
"pricing_power",
"pricing_flexibility",
"pricing_transparency",
"pricing_scalability",
"pricing_expansion",
"discount_dependency",
"premium_positioning",
"budget_positioning",
"midmarket_positioning",
"luxury_positioning",
"price_quality_ratio",
"value_for_money",
"customer_roi",
"customer_payback_period",
"price_sensitivity",
"ability_to_raise_prices",
"total_cost_of_ownership",
"switching_savings",
"premium_justification"
],

"traction": [
"revenue",
"revenue_growth",
"revenue_acceleration",
"mrr",
"mrr_growth",
"arr",
"arr_growth",
"customer_count",
"customer_growth",
"user_count",
"user_growth",
"active_users",
"dau",
"mau",
"dau_mau_ratio",
"activation_rate",
"conversion_rate",
"retention",
"cohort_retention",
"churn",
"revenue_churn",
"logo_churn",
"repeat_purchase",
"usage_frequency",
"engagement",
"organic_growth",
"viral_growth",
"referral_growth",
"word_of_mouth",
"waitlist_quality",
"waitlist_conversion",
"customer_pipeline",
"pipeline_quality",
"paid_pilots",
"enterprise_contracts",
"partnerships",
"customer_references",
"customer_testimonials",
"revenue_quality",
"revenue_predictability",
"revenue_concentration",
"traction_consistency"
],

"business_model": [
"business_model_clarity",
"business_model_simplicity",
"business_model_scalability",
"revenue_model",
"recurring_revenue",
"transaction_revenue",
"subscription_revenue",
"usage_revenue",
"marketplace_revenue",
"advertising_revenue",
"licensing_revenue",
"enterprise_revenue",
"consumer_revenue",
"revenue_diversification",
"revenue_predictability",
"revenue_expansion",
"upsell_potential",
"cross_sell_potential",
"land_and_expand",
"platform_potential"
],

"unit_economics": [
"gross_margin",
"gross_margin_potential",
"contribution_margin",
"cac",
"ltv",
"ltv_cac_ratio",
"cac_payback",
"sales_efficiency",
"magic_number",
"burn_rate",
"runway",
"capital_efficiency",
"revenue_per_employee",
"gross_profit_per_employee",
"operating_leverage",
"break_even_path",
"profitability_potential",
"margin_expansion",
"cost_reduction_potential",
"economies_of_scale"
],

"distribution": [
"distribution_strategy",
"distribution_strength",
"distribution_scalability",
"customer_acquisition",
"customer_acquisition_cost",
"sales_cycle",
"sales_cycle_predictability",
"sales_process",
"sales_repeatability",
"sales_team_quality",
"founder_led_sales",
"enterprise_sales",
"self_serve_sales",
"channel_sales",
"partnership_sales",
"marketplace_distribution",
"community_distribution",
"content_distribution",
"organic_distribution",
"paid_distribution",
"brand_distribution",
"geographic_distribution",
"distribution_moat",
"distribution_speed",
"distribution_efficiency"
],

"competition": [
"direct_competitors",
"indirect_competitors",
"substitute_products",
"incumbent_threat",
"competitive_intensity",
"competitor_growth",
"competitor_funding",
"competitor_quality",
"competitor_pricing",
"competitor_distribution",
"competitor_brand",
"competitor_customer_lockin",
"competitor_response_risk",
"competitive_position",
"market_share",
"market_share_growth",
"competitive_differentiation",
"competitive_advantage",
"competitive_durability",
"competitive_copy_risk",
"competitive_survival"
],

"switching": [
"customer_switching_motivation",
"customer_switching_cost",
"customer_switching_friction",
"customer_migration_difficulty",
"customer_learning_curve",
"customer_contract_barrier",
"customer_data_migration",
"customer_integration_barrier",
"customer_trust_barrier",
"customer_brand_barrier",
"customer_network_barrier",
"price_switching_incentive",
"quality_switching_incentive",
"speed_switching_incentive",
"convenience_switching_incentive",
"innovation_switching_incentive",
"estimated_switching_rate",
"switching_segment_size",
"switching_segment_quality",
"switching_probability"
],

"moats": [
"technology_moat",
"data_moat",
"brand_moat",
"distribution_moat",
"network_effects",
"switching_cost_moat",
"scale_moat",
"cost_moat",
"community_moat",
"regulatory_moat",
"intellectual_property_moat",
"talent_moat",
"speed_moat",
"ecosystem_moat",
"learning_moat",
"defensibility"
],

"team": [
"founder_count",
"total_team_size",
"relevant_team_size",
"team_size_stage_fit",
"team_size_problem_fit",
"team_size_capital_efficiency",
"technical_team_size",
"product_team_size",
"business_team_size",
"sales_team_size",
"marketing_team_size",
"operations_team_size",
"research_team_size",
"engineering_to_total_ratio",
"business_to_total_ratio",
"role_coverage",
"critical_role_coverage",
"skill_gap",
"team_complementarity",
"team_balance",
"team_communication",
"team_cohesion",
"team_conflict_risk",
"team_retention",
"team_hiring_velocity",
"team_hiring_quality",
"team_execution_capacity",
"team_scalability"
],

"founder": [
"founder_market_fit",
"founder_problem_understanding",
"founder_customer_understanding",
"founder_domain_expertise",
"founder_technical_expertise",
"founder_business_expertise",
"founder_product_expertise",
"founder_sales_expertise",
"founder_operational_expertise",
"founder_leadership",
"founder_vision",
"founder_communication",
"founder_learning_speed",
"founder_adaptability",
"founder_resilience",
"founder_persistence",
"founder_decision_quality",
"founder_speed",
"founder_execution",
"founder_integrity",
"founder_transparency",
"founder_recruiting_ability",
"founder_network",
"founder_credibility",
"founder_reputation",
"founder_ambition",
"founder_focus",
"founder_risk_awareness",
"founder_self_awareness",
"founder_cofounder_relationship"
],

"founder_history": [
"previous_startups",
"previous_startup_success",
"previous_startup_failure",
"previous_exit",
"previous_acquisition",
"previous_ipo_experience",
"previous_revenue_creation",
"previous_product_launches",
"previous_hiring",
"previous_fundraising",
"previous_investor_relationships",
"previous_domain_experience",
"previous_leadership",
"previous_scale_experience",
"previous_failure_learning",
"previous_execution_velocity",
"previous_achievements",
"previous_network",
"previous_reputation",
"founder_track_record"
],

"execution": [
"time_to_mvp",
"time_to_beta",
"time_to_launch",
"time_to_first_customer",
"time_to_first_revenue",
"time_to_first_100_users",
"time_to_first_1000_users",
"time_to_first_10000_users",
"time_to_first_million_revenue",
"release_frequency",
"iteration_speed",
"milestone_velocity",
"customer_growth_velocity",
"revenue_growth_velocity",
"hiring_velocity",
"partnership_velocity",
"fundraising_velocity",
"execution_consistency",
"execution_acceleration",
"execution_deceleration",
"momentum",
"momentum_acceleration",
"momentum_sustainability",
"milestone_achievement_rate",
"timeline_credibility",
"execution_predictability",
"execution_efficiency",
"speed_relative_to_competitors"
],

"technology": [
"technology_quality",
"technology_scalability",
"technology_reliability",
"technology_security",
"technology_architecture",
"technology_technical_debt",
"technology_complexity",
"technology_infrastructure",
"technology_cost",
"technology_performance",
"technology_differentiation",
"technology_replicability",
"technology_defensibility",
"ai_model_quality",
"ai_data_quality",
"ai_inference_cost",
"ai_training_cost",
"ai_accuracy",
"ai_latency",
"ai_moat"
],

"legal_regulatory": [
"regulatory_risk",
"regulatory_tailwind",
"regulatory_compliance",
"legal_risk",
"ip_ownership",
"patent_quality",
"trademark_quality",
"data_privacy",
"data_protection",
"security_compliance",
"industry_certification",
"licensing_requirements",
"regulatory_barriers",
"regulatory_moat",
"litigation_risk"
],

"capital": [
"current_valuation",
"entry_valuation",
"valuation_reasonableness",
"valuation_growth",
"cap_table_quality",
"founder_ownership",
"employee_option_pool",
"investor_quality",
"investor_support",
"previous_funding",
"funding_efficiency",
"capital_raised",
"capital_used",
"capital_remaining",
"runway",
"burn_rate",
"use_of_funds",
"future_capital_needs",
"dilution_risk",
"financing_risk"
],

"returns": [
"exit_market",
"exit_probability",
"acquisition_probability",
"ipo_probability",
"strategic_buyer_interest",
"comparable_company_quality",
"comparable_company_valuation",
"future_company_value",
"ownership_percentage",
"ownership_dilution",
"return_multiple_potential",
"return_timing",
"liquidity_potential",
"downside_protection",
"upside_potential",
"asymmetric_return",
"power_law_potential",
"expected_moic",
"expected_profit",
"probability_of_loss"
],

"evidence": [
"evidence_volume",
"evidence_quality",
"evidence_recency",
"evidence_independence",
"evidence_consistency",
"source_reliability",
"source_diversity",
"primary_source_quality",
"third_party_validation",
"claim_verifiability",
"data_completeness",
"data_contradiction",
"unknown_risk",
"information_asymmetry",
"due_diligence_confidence"
]

}

FACTORS = []

for group in FACTOR_GROUPS.values():

    FACTORS.extend(group)

FACTORS = list(dict.fromkeys(FACTORS))

print(
    "TOTAL FACTORS:",
    len(FACTORS)
)


# ============================================================
# DATABASE
# ============================================================

def init_db():

    conn = sqlite3.connect(DATABASE)

    conn.execute("""
    CREATE TABLE IF NOT EXISTS startups (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT,
        description TEXT,
        stage TEXT,
        sector TEXT,
        valuation REAL,
        created_at TEXT
    )
    """)

    conn.execute("""
    CREATE TABLE IF NOT EXISTS founders (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        startup_id INTEGER,
        name TEXT,
        linkedin TEXT,
        personal_website TEXT,
        portfolio TEXT,
        github TEXT,
        other_links TEXT
    )
    """)

    conn.execute("""
    CREATE TABLE IF NOT EXISTS sources (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        startup_id INTEGER,
        url TEXT
    )
    """)

    conn.execute("""
    CREATE TABLE IF NOT EXISTS documents (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        startup_id INTEGER,
        filename TEXT,
        path TEXT
    )
    """)

    conn.execute("""
    CREATE TABLE IF NOT EXISTS analyses (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        startup_id INTEGER,
        analysis TEXT,
        created_at TEXT
    )
    """)

    conn.commit()
    conn.close()


init_db()


# ============================================================
# DATA MODELS
# ============================================================

class Founder(BaseModel):

    name: str

    linkedin: Optional[str] = None

    personal_website: Optional[str] = None

    portfolio: Optional[str] = None

    github: Optional[str] = None

    other_links: List[str] = Field(
        default_factory=list
    )


class Startup(BaseModel):

    name: str

    description: Optional[str] = None

    stage: str = "Unknown"

    sector: str = "Unknown"

    valuation: Optional[float] = None

    founders: List[Founder] = Field(
        default_factory=list
    )

    sources: List[str] = Field(
        default_factory=list
    )


class AnalyzeRequest(BaseModel):

    startup_id: int


# ============================================================
# URL CRAWLER
# ============================================================

def normalize_url(url):

    url = url.strip()

    if not url:
        return ""

    if not url.startswith(
        ("http://", "https://")
    ):
        url = "https://" + url

    return url


def crawl_url(url):

    url = normalize_url(url)

    try:

        response = requests.get(
            url,
            timeout=20,
            headers={
                "User-Agent":
                "VC-Brain-Research-Agent/5.0"
            }
        )

        soup = BeautifulSoup(
            response.text,
            "html.parser"
        )

        title = ""

        if soup.title:

            title = soup.title.get_text(
                " ",
                strip=True
            )

        for tag in soup(
            [
                "script",
                "style",
                "noscript",
                "svg"
            ]
        ):

            tag.extract()

        text = soup.get_text(
            " ",
            strip=True
        )

        text = re.sub(
            r"\s+",
            " ",
            text
        )

        return {

            "source_type":
            "website",

            "source":
            url,

            "title":
            title,

            "text":
            text[:MAX_URL_TEXT],

            "status":
            "FETCHED"

        }

    except Exception as e:

        return {

            "source_type":
            "website",

            "source":
            url,

            "title":
            "",

            "text":
            "",

            "status":
            "UNVERIFIED",

            "error":
            str(e)

        }


# ============================================================
# DOCUMENT EXTRACTION
# ============================================================

def extract_pdf(path):

    text = ""

    try:

        reader = PdfReader(path)

        for page in reader.pages:

            text += (
                page.extract_text()
                or ""
            )

    except Exception as e:

        text = (
            "PDF EXTRACTION ERROR: "
            + str(e)
        )

    return text


def extract_docx(path):

    text = ""

    try:

        document = Document(path)

        for paragraph in document.paragraphs:

            text += (
                paragraph.text
                + "\n"
            )

        for table in document.tables:

            for row in table.rows:

                text += " | ".join(
                    cell.text
                    for cell in row.cells
                )

                text += "\n"

    except Exception as e:

        text = (
            "DOCX EXTRACTION ERROR: "
            + str(e)
        )

    return text


def extract_pptx(path):

    text = ""

    try:

        presentation = Presentation(path)

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1
        ):

            text += (
                f"\nSLIDE {slide_number}\n"
            )

            for shape in slide.shapes:

                if hasattr(
                    shape,
                    "text"
                ):

                    text += (
                        shape.text
                        + "\n"
                    )

    except Exception as e:

        text = (
            "PPTX EXTRACTION ERROR: "
            + str(e)
        )

    return text


def extract_txt(path):

    try:

        with open(
            path,
            "r",
            encoding="utf-8",
            errors="ignore"
        ) as file:

            return file.read()

    except Exception as e:

        return (
            "TXT EXTRACTION ERROR: "
            + str(e)
        )


def extract_document(path):

    extension = path.lower().split(".")[-1]

    if extension == "pdf":

        return extract_pdf(path)

    if extension == "docx":

        return extract_docx(path)

    if extension == "pptx":

        return extract_pptx(path)

    if extension in [
        "txt",
        "md",
        "csv"
    ]:

        return extract_txt(path)

    return (
        "UNSUPPORTED DOCUMENT TYPE"
    )


# ============================================================
# AI JSON
# ============================================================

def ask_ai(system_prompt, user_prompt):

    try:

        system_prompt = (
            system_prompt
            + "\n\nReturn the response as valid JSON."
        )

        response = client.chat.completions.create(

            model="gpt-4o",

            temperature=0.1,

            response_format={
                "type": "json_object"
            },

            messages=[

                {
                    "role": "system",
                    "content": system_prompt
                },

                {
                    "role": "user",
                    "content": user_prompt
                }

            ]

        )

        content = response.choices[0].message.content

        if not content:

            raise RuntimeError(
                "OpenAI returned an empty response"
            )

        return json.loads(content)

    except Exception as e:

        print(
            "OPENAI ERROR:",
            repr(e)
        )

        raise HTTPException(
            status_code=500,
            detail=f"AI analysis failed: {str(e)}"
        )


# ============================================================
# EVIDENCE EXTRACTION
# ============================================================

def extract_evidence(
    startup,
    evidence
):

    return ask_ai(

"""
You are an evidence extraction engine for a venture capital
investment decision.

Analyze EVERY source supplied.

Sources may include:

- company websites
- product websites
- founder LinkedIn pages
- personal websites
- GitHub
- portfolios
- press
- competitor websites
- market reports
- pitch decks
- PDFs
- DOCX files
- PPTX files
- financial documents
- milestone documents
- research documents

For every meaningful claim:

- claim
- evidence
- source
- source_type
- date if available
- verification_status
- confidence

Verification statuses MUST be:

VERIFIED
PARTIALLY_VERIFIED
UNVERIFIED
CONTRADICTED

Rules:

1. Never invent information.
2. A company claim is not automatically independently verified.
3. If multiple sources support a claim, increase confidence.
4. If sources conflict, mark CONTRADICTED.
5. If evidence is incomplete, mark PARTIALLY_VERIFIED.
6. Missing evidence should become an important unknown.
7. Extract dated milestones.
8. Extract founder history.
9. Extract competitor information.
10. Extract pricing and quality information.
11. Extract traction and financial claims.

Return:

{
  "claims": [],
  "milestones": [],
  "contradictions": [],
  "important_unknowns": [],
  "source_quality": []
}
""",

        json.dumps(
            {
                "startup":
                startup,

                "evidence":
                evidence
            },

            indent=2
        )
    )


# ============================================================
# FOUNDER + TEAM
# ============================================================

def analyze_team(
    startup,
    founders,
    evidence
):

    return ask_ai(

"""
You are a senior venture capital team diligence analyst.

Analyze EACH founder separately and then the team collectively.

For every founder analyze evidence for:

- identity
- role
- domain expertise
- technical expertise
- business expertise
- product expertise
- sales expertise
- operations expertise
- founder-market fit
- leadership
- communication
- learning speed
- adaptability
- resilience
- persistence
- decision quality
- execution
- integrity
- transparency
- recruiting ability
- network
- credibility
- ambition
- focus
- previous startups
- successful startups
- failed startups
- exits
- acquisitions
- previous revenue
- previous launches
- previous hiring
- previous fundraising
- previous leadership
- previous domain experience
- previous execution velocity
- achievements

Analyze the entire team:

- founder count
- total team size
- relevant team size
- technical team
- product team
- business team
- sales team
- operations team
- research team
- role coverage
- critical gaps
- complementarity
- team balance
- cohesion
- conflict risk
- hiring ability
- execution capacity
- team size relative to stage
- team size relative to problem complexity
- capital efficiency

A failed startup is NOT automatically negative.

A large team is NOT automatically positive.

A small team is NOT automatically negative.

Return evidence-backed analysis.
""",

        json.dumps(
            {
                "startup":
                startup,

                "founders":
                founders,

                "evidence":
                evidence
            },

            indent=2
        )
    )


# ============================================================
# EXECUTION TIMELINE
# ============================================================

def analyze_execution(
    startup,
    evidence
):

    return ask_ai(

"""
You are an execution velocity analyst.

Reconstruct the startup timeline from actual evidence.

Extract dates for:

- founding
- MVP
- beta
- public launch
- first customer
- first revenue
- first 100 users
- first 1,000 users
- first 10,000 users
- revenue milestones
- product releases
- hires
- partnerships
- fundraising

Calculate where possible:

- time to MVP
- time to beta
- time to launch
- time to first customer
- time to first revenue
- customer growth velocity
- revenue growth velocity
- product release velocity
- hiring velocity
- partnership velocity
- milestone velocity
- acceleration
- deceleration
- consistency
- momentum
- execution efficiency
- speed relative to competitors

Never invent dates.

Return:

{
  "timeline": [],
  "velocity_metrics": {},
  "momentum": {},
  "execution_strengths": [],
  "execution_risks": [],
  "unknowns": []
}
""",

        json.dumps(
            {
                "startup":
                startup,

                "evidence":
                evidence
            },

            indent=2
        )
    )


# ============================================================
# COMPETITION
# ============================================================

def analyze_competition(
    startup,
    evidence
):

    return ask_ai(

"""
You are a competitive intelligence and market-switching analyst.

Identify:

- direct competitors
- indirect competitors
- substitutes
- incumbents
- current customer alternatives

Compare:

- price
- quality
- features
- speed
- convenience
- distribution
- brand
- trust
- switching cost
- customer lock-in

Analyze:

1. What is genuinely different?
2. Is the difference valuable?
3. Is it difficult to copy?
4. Is the price lower?
5. Is the quality higher?
6. Is the product faster?
7. Is the product more convenient?
8. Will customers realistically switch?
9. Which segments switch first?
10. Which segments will not switch?
11. Does the product create a new category?
12. Can competitors easily respond?

Estimate switching probability as a RANGE:

low
base
high

Do not present the estimate as certainty.

Return:

{
  "competitors": [],
  "substitutes": [],
  "comparison": [],
  "differentiation": {},
  "switching_analysis": {},
  "competitive_risks": [],
  "competitive_opportunities": []
}
""",

        json.dumps(
            {
                "startup":
                startup,

                "evidence":
                evidence
            },

            indent=2
        )
    )


# ============================================================
# FACTOR SCORING
# ============================================================

def score_all_factors(
    startup,
    evidence,
    team,
    execution,
    competition
):

    all_factor_scores = {}

    # Score factors in batches instead of asking the AI to return
    # all 542 factors in one enormous response.
    BATCH_SIZE = 40

    for i in range(0, len(FACTORS), BATCH_SIZE):

        factor_batch = FACTORS[i:i + BATCH_SIZE]

        print(
            f"Scoring factors {i + 1} "
            f"to {min(i + BATCH_SIZE, len(FACTORS))} "
            f"of {len(FACTORS)}..."
        )

        system_prompt = f"""
You are the core VC scoring engine.

You are scoring a batch of factors from a larger VC investment analysis.

Score EVERY factor in this batch.

FACTORS IN THIS BATCH:

{json.dumps(factor_batch, indent=2)}

For EACH factor return:

{{
  "score": 0,
  "status": "VERIFIED",
  "confidence": 0,
  "evidence": [],
  "reasoning": "",
  "unknowns": []
}}

Allowed status values:

VERIFIED
PARTIALLY_VERIFIED
UNVERIFIED
CONTRADICTED

Rules:

- Never invent facts.
- Use only the supplied evidence.
- If evidence is absent, mark UNVERIFIED.
- If evidence is incomplete, mark PARTIALLY_VERIFIED.
- If sources conflict, mark CONTRADICTED.
- A failed startup is not automatically negative.
- Judge team size relative to stage and business complexity.
- Analyze founders individually where relevant.
- Consider previous companies, including failures and exits.
- Consider execution velocity.
- Consider milestone timing.
- Consider team composition.
- Consider competitors and substitutes.
- Consider customer switching probability.
- Consider price versus quality.
- Consider market growth.
- Consider innovation.
- Consider defensibility.
- Consider financial return potential.

Return valid JSON only in this exact structure:

{{
  "factor_scores": {{
    "factor_name": {{
      "score": 0,
      "status": "UNVERIFIED",
      "confidence": 0,
      "evidence": [],
      "reasoning": "",
      "unknowns": []
    }}
  }}
}}

Every factor in this batch MUST be present exactly once.
"""

        user_prompt = json.dumps(
            {
                "startup": startup,
                "evidence": evidence,
                "team": team,
                "execution": execution,
                "competition": competition
            },
            indent=2
        )

        batch_result = ask_ai(
            system_prompt,
            user_prompt
        )

        batch_scores = batch_result.get(
            "factor_scores",
            {}
        )

        all_factor_scores.update(
            batch_scores
        )

    print(
        f"Completed scoring {len(all_factor_scores)} "
        f"of {len(FACTORS)} factors."
    )

    return {
        "factor_scores": all_factor_scores
    }


# ============================================================
# CATEGORY SCORING
# ============================================================

def category_scores(
    factor_data
):

    output = {}

    for category, factors in FACTOR_GROUPS.items():

        values = []

        verified_values = []

        for factor in factors:

            item = factor_data.get(
                factor
            )

            if not item:

                continue

            score = item.get(
                "score"
            )

            if isinstance(
                score,
                (int, float)
            ):

                values.append(
                    score
                )

                if item.get(
                    "status"
                ) == "VERIFIED":

                    verified_values.append(
                        score
                    )

        output[category] = {

            "score":

            round(
                sum(values) / len(values),
                2
            )
            if values
            else 0,

            "evidence_confidence":

            round(
                len(verified_values)
                / len(values)
                * 100,
                2
            )
            if values
            else 0,

            "factor_count":

            len(values)

        }

    return output


# ============================================================
# STAGE WEIGHTS
# ============================================================

def stage_weights(stage):

    stage = stage.lower()

    weights = {
        category: 1.0
        for category in FACTOR_GROUPS
    }

    if (
        "pre" in stage
        or "seed" in stage
    ):

        weights.update({

            "problem": 1.15,

            "customer": 1.10,

            "market": 1.20,

            "product": 1.10,

            "innovation": 1.05,

            "traction": 0.80,

            "unit_economics": 0.80,

            "team": 1.20,

            "founder": 1.20,

            "founder_history": 1.05,

            "execution": 1.20,

            "evidence": 1.15

        })

    elif "series a" in stage:

        weights.update({

            "market": 1.10,

            "traction": 1.20,

            "unit_economics": 1.15,

            "team": 1.05,

            "execution": 1.10,

            "competition": 1.10,

            "moats": 1.10,

            "evidence": 1.10

        })

    return weights


def overall_score(
    categories,
    stage
):

    weights = stage_weights(
        stage
    )

    total = 0

    weight_total = 0

    for category, data in categories.items():

        weight = weights.get(
            category,
            1
        )

        total += (
            data["score"]
            * weight
        )

        weight_total += weight

    return round(
        total / weight_total,
        2
    )


# ============================================================
# $100K RETURN MODEL
# ============================================================

def return_model(
    startup,
    categories,
    score,
    factor_scores,
    team,
    execution,
    competition
):

    return ask_ai(

"""
You are a venture capital return-modeling analyst.

The proposed investment is exactly $100,000.

Analyze:

- market size
- growth
- founder quality
- team
- relevant team size
- execution velocity
- momentum
- product
- innovation
- competition
- customer switching
- traction
- economics
- valuation
- dilution
- exit potential
- evidence quality
- unknowns

Estimate probabilities for:

1. Loss / severe loss
2. Approximately returning capital
3. 1.5x to 3x
4. 3x to 10x
5. 10x+

Probabilities MUST sum to 100.

Calculate:

expected MOIC
expected value of $100,000
expected profit/loss

Use:

expected value =
sum(
probability × outcome value
)

Do not pretend these are precise predictions.

Return:

{
  "investment_amount": 100000,
  "probabilities": {},
  "scenario_values": {},
  "expected_moic": 0,
  "expected_value": 0,
  "expected_profit": 0,
  "downside_case": "",
  "base_case": "",
  "upside_case": "",
  "confidence": 0,
  "reasoning": ""
}
""",

        json.dumps(
            {
                "startup":
                startup,

                "categories":
                categories,

                "overall_score":
                score,

                "factor_scores":
                factor_scores,

                "team":
                team,

                "execution":
                execution,

                "competition":
                competition
            },

            indent=2
        )
    )


# ============================================================
# FINAL DECISION
# ============================================================

def final_decision(
    startup,
    score,
    categories,
    returns,
    team,
    execution,
    competition,
    factors
):

    return ask_ai(

"""
You are the final investment committee.

The question is:

Should the VC invest $100,000?

Choose exactly one:

INVEST
INVESTIGATE
PASS

INVEST:

Evidence and risk-adjusted return justify deploying $100,000.

INVESTIGATE:

Potential is strong enough to continue diligence,
but material unknowns or contradictions must be resolved.

PASS:

Risk-adjusted return is not attractive enough.

Consider:

- all factor scores
- verification status
- category scores
- stage
- sector
- founders
- team size
- relevant team size
- team composition
- founder history
- founder-market fit
- execution velocity
- momentum
- market
- product
- innovation
- competitors
- differentiation
- switching probability
- traction
- unit economics
- valuation
- return model
- probability of loss
- expected MOIC
- expected profit
- evidence confidence
- contradictions
- unknowns

Return:

{
  "decision": "INVEST",
  "investment_amount": 100000,
  "confidence": 0,
  "investment_thesis": "",
  "reasoning": "",
  "probability_of_loss": 0,
  "probability_of_profit": 0,
  "expected_moic": 0,
  "expected_value": 0,
  "expected_profit": 0,
  "top_reasons_to_invest": [],
  "top_reasons_not_to_invest": [],
  "key_risks": [],
  "critical_diligence_questions": []
}
""",

        json.dumps(
            {
                "startup":
                startup,

                "overall_score":
                score,

                "categories":
                categories,

                "returns":
                returns,

                "team":
                team,

                "execution":
                execution,

                "competition":
                competition,

                "factors":
                factors
            },

            indent=2
        )
    )


# ============================================================
# CREATE STARTUP
# ============================================================

@app.post("/startups")
def create_startup(
    startup: Startup
):

    conn = sqlite3.connect(
        DATABASE
    )

    cursor = conn.execute(

"""
INSERT INTO startups
(
name,
description,
stage,
sector,
valuation,
created_at
)
VALUES (?, ?, ?, ?, ?, ?)
""",

        (
            startup.name,
            startup.description,
            startup.stage,
            startup.sector,
            startup.valuation,
            datetime.utcnow().isoformat()
        )
    )

    startup_id = cursor.lastrowid

    for founder in startup.founders:

        conn.execute(

"""
INSERT INTO founders
(
startup_id,
name,
linkedin,
personal_website,
portfolio,
github,
other_links
)
VALUES (?, ?, ?, ?, ?, ?, ?)
""",

            (
                startup_id,
                founder.name,
                founder.linkedin,
                founder.personal_website,
                founder.portfolio,
                founder.github,
                json.dumps(
                    founder.other_links
                )
            )
        )

        links = [

            founder.linkedin,

            founder.personal_website,

            founder.portfolio,

            founder.github

        ]

        links.extend(
            founder.other_links
        )

        for link in links:

            if link:

                conn.execute(

"""
INSERT INTO sources
(
startup_id,
url
)
VALUES (?, ?)
""",

                    (
                        startup_id,
                        link
                    )
                )

    for source in startup.sources:

        if source:

            conn.execute(

"""
INSERT INTO sources
(
startup_id,
url
)
VALUES (?, ?)
""",

                (
                    startup_id,
                    source
                )
            )

    conn.commit()
    conn.close()

    return {
        "id":
        startup_id,

        "message":
        "Startup created"
    }


# ============================================================
# UPLOAD DOCUMENT
# ============================================================

@app.post(
    "/startups/{startup_id}/documents"
)
async def upload_document(
    startup_id: int,
    file: UploadFile = File(...)
):

    safe_name = re.sub(
        r"[^a-zA-Z0-9._-]",
        "_",
        file.filename
    )

    path = os.path.join(
        UPLOAD_DIR,
        f"{startup_id}_{safe_name}"
    )

    content = await file.read()

    with open(
        path,
        "wb"
    ) as output:

        output.write(
            content
        )

    conn = sqlite3.connect(
        DATABASE
    )

    conn.execute(

"""
INSERT INTO documents
(
startup_id,
filename,
path
)
VALUES (?, ?, ?)
""",

        (
            startup_id,
            safe_name,
            path
        )
    )

    conn.commit()
    conn.close()

    return {

        "uploaded":
        True,

        "filename":
        safe_name,

        "text_extraction":
        "will occur during analysis"

    }


# ============================================================
# FULL ANALYSIS PIPELINE
# ============================================================

@app.post("/analyze")
def analyze(
    request: AnalyzeRequest
):

    conn = sqlite3.connect(
        DATABASE
    )

    conn.row_factory = sqlite3.Row

    startup = conn.execute(

"""
SELECT *
FROM startups
WHERE id = ?
""",

        (
            request.startup_id,
        )
    ).fetchone()

    founders = conn.execute(

"""
SELECT *
FROM founders
WHERE startup_id = ?
""",

        (
            request.startup_id,
        )
    ).fetchall()

    sources = conn.execute(

"""
SELECT *
FROM sources
WHERE startup_id = ?
""",

        (
            request.startup_id,
        )
    ).fetchall()

    documents = conn.execute(

"""
SELECT *
FROM documents
WHERE startup_id = ?
""",

        (
            request.startup_id,
        )
    ).fetchall()

    conn.close()

    if not startup:

        raise HTTPException(
            404,
            "Startup not found"
        )

    startup = dict(
        startup
    )

    founders = [
        dict(x)
        for x in founders
    ]

    # --------------------------------------------------------
    # 1. CRAWL ALL URLS
    # --------------------------------------------------------

    evidence = []

    seen = set()

    for source in sources:

        url = source["url"]

        if not url:

            continue

        normalized = normalize_url(
            url
        )

        if normalized in seen:

            continue

        seen.add(
            normalized
        )

        page = crawl_url(
            normalized
        )

        evidence.append(
            page
        )

    # --------------------------------------------------------
    # 2. EXTRACT ALL UPLOADED DOCUMENTS
    # --------------------------------------------------------

    for document in documents:

        text = extract_document(
            document["path"]
        )

        evidence.append({

            "source_type":
            "uploaded_document",

            "source":
            document["filename"],

            "title":
            document["filename"],

            "text":
            text[
                :MAX_DOCUMENT_TEXT
            ],

            "status":
            "EXTRACTED"

        })

    # --------------------------------------------------------
    # 3. LIMIT TOTAL EVIDENCE
    # --------------------------------------------------------

    total = 0

    combined_evidence = []

    for item in evidence:

        text = item.get(
            "text",
            ""
        )

        if (
            total
            + len(text)
            > MAX_TOTAL_EVIDENCE
        ):

            remaining = (
                MAX_TOTAL_EVIDENCE
                - total
            )

            item["text"] = text[
                :max(
                    0,
                    remaining
                )
            ]

        combined_evidence.append(
            item
        )

        total += len(
            item.get(
                "text",
                ""
            )
        )

        if total >= MAX_TOTAL_EVIDENCE:

            break

    # --------------------------------------------------------
    # 4. EXTRACT EVIDENCE
    # --------------------------------------------------------

    extracted = extract_evidence(
        startup,
        combined_evidence
    )

    # --------------------------------------------------------
    # 5. ANALYZE FOUNDERS AND TEAM
    # --------------------------------------------------------

    team = analyze_team(
        startup,
        founders,
        extracted
    )

    # --------------------------------------------------------
    # 6. RECONSTRUCT EXECUTION
    # --------------------------------------------------------

    execution = analyze_execution(
        startup,
        extracted
    )

    # --------------------------------------------------------
    # 7. COMPETITION
    # --------------------------------------------------------

    competition = analyze_competition(
        startup,
        extracted
    )

    # --------------------------------------------------------
    # 8. SCORE ALL FACTORS
    # --------------------------------------------------------

    factor_output = score_all_factors(
        startup,
        extracted,
        team,
        execution,
        competition
    )

    factor_scores = factor_output.get(
        "factor_scores",
        {}
    )

    # --------------------------------------------------------
    # 9. CATEGORY SCORES
    # --------------------------------------------------------

    categories = category_scores(
        factor_scores
    )

    # --------------------------------------------------------
    # 10. STAGE-ADJUSTED SCORE
    # --------------------------------------------------------

    score = overall_score(
        categories,
        startup["stage"]
    )

    # --------------------------------------------------------
    # 11. RETURN MODEL
    # --------------------------------------------------------

    returns = return_model(
        startup,
        categories,
        score,
        factor_scores,
        team,
        execution,
        competition
    )

    # --------------------------------------------------------
    # 12. FINAL DECISION
    # --------------------------------------------------------

    decision = final_decision(
        startup,
        score,
        categories,
        returns,
        team,
        execution,
        competition,
        factor_scores
    )

    result = {

        "startup":
        startup,

        "factor_count":
        len(FACTORS),

        "sources_crawled":
        len(
            [
                x
                for x in evidence
                if x["source_type"]
                == "website"
            ]
        ),

        "documents_extracted":
        len(
            [
                x
                for x in evidence
                if x["source_type"]
                == "uploaded_document"
            ]
        ),

        "evidence":
        extracted,

        "team_analysis":
        team,

        "execution_analysis":
        execution,

        "competition_analysis":
        competition,

        "factor_scores":
        factor_scores,

        "category_scores":
        categories,

        "overall_score":
        score,

        "return_model":
        returns,

        "investment_decision":
        decision,

        "created_at":
        datetime.utcnow().isoformat()

    }

    conn = sqlite3.connect(
        DATABASE
    )

    conn.execute(

"""
INSERT INTO analyses
(
startup_id,
analysis,
created_at
)
VALUES (?, ?, ?)
""",

        (
            request.startup_id,

            json.dumps(
                result
            ),

            datetime.utcnow().isoformat()
        )
    )

    conn.commit()
    conn.close()

    return result


# ============================================================
# GET ANALYSIS
# ============================================================

@app.get(
    "/startups/{startup_id}/analysis"
)
def get_analysis(
    startup_id: int
):

    conn = sqlite3.connect(
        DATABASE
    )

    conn.row_factory = sqlite3.Row

    row = conn.execute(

"""
SELECT analysis
FROM analyses
WHERE startup_id = ?
ORDER BY id DESC
LIMIT 1
""",

        (
            startup_id,
        )
    ).fetchone()

    conn.close()

    if not row:

        raise HTTPException(
            404,
            "No analysis found"
        )

    return json.loads(
        row["analysis"]
    )


# ============================================================
# FRONTEND
# ============================================================

app.mount(
    "/frontend",
    StaticFiles(
        directory="frontend"
    ),
    name="frontend"
)


@app.get("/")
def root():

    return {

        "system":
        "VC Brain",

        "factor_count":
        len(FACTORS),

        "investment_amount":
        100000,

        "status":
        "operational"

    }


@app.get("/app")
def frontend():

    return FileResponse(
        "frontend/index.html"
    )